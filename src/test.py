from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Slide 1: About Yourself
slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])
title1 = slide1.shapes.title
title1.text = "About Me"

content1 = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
frame1 = content1.text_frame
frame1.text = (
    "Ali Ahmadi\n"
    "Rasht, Iran\n"
    "Email: ali.ahmadi.k79@gmail.com | Phone: +98 903-312-2974\n"
    "LinkedIn: linkedin.com/in/ali-ahmadi-7023721ba/\n\n"
    "I am passionate about merging AI technologies with cultural heritage and music.\n"
    "With an MSc in Software Engineering and research experience in audio processing,\n"
    "I aim to develop innovative solutions to enhance multimodal AI systems."
)

# Slide 2: Contribution to a Project or Paper
slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
title2 = slide2.shapes.title
title2.text = "My Contribution to Research"

content2 = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
frame2 = content2.text_frame
frame2.text = (
    "Ensemble of Contrastive Models for Instrument Classification\n"
    "\n"
    "- Developed models to classify instruments in classical Persian music using\n"
    "  a comprehensive dataset.\n"
    "- Addressed challenges unique to Iranian music through deep learning\n"
    "  and time-series data analysis.\n"
    "- Enhanced understanding of audio processing and machine learning.\n"
    "\n"
    "Outcome: Improved classification accuracy and a deeper insight into\n"
    "Persian musical instruments. The research is under review for publication."
)

# Slide 3: Contribution to HAICu Project
slide3 = presentation.slides.add_slide(presentation.slide_layouts[5])
title3 = slide3.shapes.title
title3.text = "My Ideas for HAICu Project"

content3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
frame3 = content3.text_frame
frame3.text = (
    "How I Can Contribute:\n"
    "\n"
    "- Multimodal AI: Utilize my expertise in audio processing, machine learning,\n"
    "  and NLP to develop advanced models for analyzing speech and music.\n"
    "- Data-Driven Journalism: Create tools that enable cultural and\n"
    "  journalistic exploration using multimodal archives.\n"
    "- Collaboration: Work with HAICu teams to design AI systems that bridge\n"
    "  cultural heritage and modern technology.\n"
    "\n"
    "Vision: Build AI-driven solutions for preserving and unlocking insights\n"
    "from cultural digital archives, promoting public engagement."
)

# Save the presentation
presentation.save("HAICu_Interview_Slides.pptx")
